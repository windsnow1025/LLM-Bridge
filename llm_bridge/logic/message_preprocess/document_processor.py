import logging
from io import BytesIO

import fitz
import openpyxl
from docxlatex import Document as DocxLatexDocument
from fastapi import HTTPException
from pptx import Presentation

from llm_bridge.logic.file_fetch import fetch_file_data
from llm_bridge.logic.message_preprocess import file_type_checker


async def extract_text_from_file(file_url: str) -> str:
    file_content, content_type = await fetch_file_data(file_url)

    if file_type_checker.is_file_type_supported(file_url) is False:
        raise HTTPException(status_code=415,
                            detail=f"legacy filetypes ('.doc', '.xls', '.ppt') are not supported - {file_url}")

    file_type, sub_type = await file_type_checker.get_file_type(file_url)

    try:
        if sub_type == "code":
            return extract_text_from_code(file_content)
        if sub_type == "pdf":
            return extract_text_from_pdf(file_content)
        if sub_type == "word":
            return extract_text_from_word(file_content)
        if sub_type == "excel":
            return extract_text_from_excel(file_content)
        if sub_type == "powerpoint":
            return extract_text_from_ppt(file_content)
    except Exception as e:
        logging.exception(e)
        raise HTTPException(
            status_code=415,
            detail=f"Error during text extraction - {file_url} ({file_type}/{sub_type}) - {str(e)}"
        )

    try:
        return file_content.decode('utf-8')
    except UnicodeDecodeError as e:
        message = f"Unicode Decode Error during text extraction - {file_url} ({file_type}/{sub_type}) - {str(e)}"
        logging.error(message)
        return message


def extract_text_from_code(file_content: bytes) -> str:
    return file_content.decode('utf-8')


def extract_text_from_pdf(file_content: bytes) -> str:
    text = ""
    with fitz.open(stream=file_content, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text


def extract_text_from_word(file_content: bytes) -> str:
    with BytesIO(file_content) as f:
        doc = DocxLatexDocument(f)
        text_with_equations = doc.get_text()
        return text_with_equations


def extract_text_from_excel(file_content: bytes) -> str:
    text = ""
    workbook = openpyxl.load_workbook(BytesIO(file_content), data_only=True)
    for sheet in workbook:
        text += f"Sheet: {sheet.title}\n"
        for row in sheet.iter_rows(values_only=True):
            text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return text


def extract_text_from_ppt(file_content: bytes) -> str:
    text = ""
    presentation = Presentation(BytesIO(file_content))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
